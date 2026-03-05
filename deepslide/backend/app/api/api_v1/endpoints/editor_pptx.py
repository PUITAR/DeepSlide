import os
import re
from typing import Any, Dict, List

try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover
    fitz = None


def _build_pptx_from_images(image_paths: List[str], out_path: str):
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image

    prs = Presentation()
    prs.slide_width = Emu(13.333 * 914400)
    prs.slide_height = Emu(7.5 * 914400)
    blank = prs.slide_layouts[6]

    sw = float(prs.slide_width)
    sh = float(prs.slide_height)

    for p in image_paths:
        slide = prs.slides.add_slide(blank)
        try:
            with Image.open(p) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = (1920, 1080)

        scale = min(sw / float(iw), sh / float(ih))
        w = int(float(iw) * scale)
        h = int(float(ih) * scale)
        left = int((sw - w) / 2)
        top = int((sh - h) / 2)
        slide.shapes.add_picture(p, left, top, width=w, height=h)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)


def _build_pptx_from_pdf_editable(pdf_path: str, out_path: str) -> bool:
    if not pdf_path or not os.path.exists(pdf_path):
        return False
    if fitz is None:
        return False

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception:
        return False

    import tempfile

    def _standard_slide_size_from_ratio(r: float):
        if abs(r - 4 / 3) < 0.05:
            return (10.0, 7.5)
        if abs(r - 16 / 9) < 0.05:
            return (13.333, 7.5)
        h = 7.5
        return (h * r, h)

    def _rgb_from_int(c: int):
        r = (c >> 16) & 255
        g = (c >> 8) & 255
        b = c & 255
        return r, g, b

    def _is_bold(span: Dict[str, Any]) -> bool:
        f = str(span.get("font", "") or "")
        return ("bold" in f.lower()) or ("black" in f.lower())

    def _is_italic(span: Dict[str, Any]) -> bool:
        f = str(span.get("font", "") or "")
        fl = f.lower()
        return ("italic" in fl) or ("oblique" in fl)

    try:
        doc = fitz.open(pdf_path)
    except Exception:
        return False

    try:
        if doc.page_count <= 0:
            return False
        w_pt, h_pt = float(doc[0].rect.width), float(doc[0].rect.height)
        ratio = (w_pt / h_pt) if h_pt else (4 / 3)
        slide_w_in, slide_h_in = _standard_slide_size_from_ratio(ratio)

        prs = Presentation()
        prs.slide_width = Inches(slide_w_in)
        prs.slide_height = Inches(slide_h_in)
        blank = prs.slide_layouts[6]

        slide_w_pt = slide_w_in * 72.0
        slide_h_pt = slide_h_in * 72.0
        sx = slide_w_pt / w_pt if w_pt else 1.0
        sy = slide_h_pt / h_pt if h_pt else 1.0
        font_scale = sy

        tmp_dir = tempfile.mkdtemp(prefix="pdf2pptx_")

        for page_idx in range(doc.page_count):
            page = doc.load_page(int(page_idx))
            d = page.get_text("dict")
            slide = prs.slides.add_slide(blank)

            for block in d.get("blocks", []) or []:
                if block.get("type") != 1:
                    continue
                bbox = block.get("bbox")
                img_bytes = block.get("image")
                if not bbox or not isinstance(img_bytes, (bytes, bytearray)):
                    continue
                try:
                    x0, y0, x1, y1 = bbox
                except Exception:
                    continue
                left = Inches((float(x0) * sx) / 72.0)
                top = Inches((float(y0) * sy) / 72.0)
                width = Inches(((float(x1) - float(x0)) * sx) / 72.0)
                height = Inches(((float(y1) - float(y0)) * sy) / 72.0)
                ext = str(block.get("ext") or "png").strip().lower()
                if not re.fullmatch(r"[a-z0-9]{2,6}", ext or ""):
                    ext = "png"
                img_path = os.path.join(tmp_dir, f"p{page_idx+1}_img{int(block.get('number', 0) or 0)}.{ext}")
                try:
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                except Exception:
                    pass

            text_blocks = [b for b in (d.get("blocks", []) or []) if b.get("type") == 0]
            text_blocks.sort(key=lambda b: ((b.get("bbox") or [0, 0, 0, 0])[1], (b.get("bbox") or [0, 0, 0, 0])[0]))

            for block in text_blocks:
                bbox = block.get("bbox")
                if not bbox:
                    continue
                try:
                    x0, y0, x1, y1 = bbox
                except Exception:
                    continue
                if (float(x1) - float(x0)) < 2 or (float(y1) - float(y0)) < 2:
                    continue

                left = Inches((float(x0) * sx) / 72.0)
                top = Inches((float(y0) * sy) / 72.0)
                width = Inches(((float(x1) - float(x0)) * sx) / 72.0)
                height = Inches(((float(y1) - float(y0)) * sy) / 72.0)

                shape = slide.shapes.add_textbox(left, top, width, height)
                tf = shape.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.margin_left = 0
                tf.margin_right = 0
                tf.margin_top = 0
                tf.margin_bottom = 0

                first_par = True
                for line in block.get("lines", []) or []:
                    p = tf.paragraphs[0] if first_par else tf.add_paragraph()
                    first_par = False
                    try:
                        p.space_before = Pt(0)
                        p.space_after = Pt(0)
                        p.line_spacing = 1.0
                    except Exception:
                        pass

                    for span in line.get("spans", []) or []:
                        txt = str(span.get("text") or "")
                        if not txt:
                            continue
                        txt = txt.replace("\u00ad", "")
                        run = p.add_run()
                        run.text = txt
                        font = run.font

                        try:
                            size_pt = float(span.get("size") or 12.0)
                            font.size = Pt(max(1.0, size_pt * float(font_scale)))
                        except Exception:
                            pass

                        try:
                            font_name = span.get("font")
                            if font_name:
                                font.name = str(font_name)
                        except Exception:
                            pass

                        try:
                            font.bold = _is_bold(span)
                            font.italic = _is_italic(span)
                        except Exception:
                            pass

                        c = span.get("color")
                        if isinstance(c, int):
                            try:
                                r, g, b = _rgb_from_int(int(c))
                                font.color.rgb = RGBColor(r, g, b)
                            except Exception:
                                pass

        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        prs.save(out_path)
        return True
    except Exception:
        return False
    finally:
        try:
            doc.close()
        except Exception:
            pass


def _build_pptx_from_tex(recipe_dir: str, out_path: str) -> bool:
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except Exception:
        return False

    content_path = os.path.join(recipe_dir, "content.tex")
    if not os.path.exists(content_path):
        return False

    try:
        with open(content_path, "r", encoding="utf-8", errors="ignore") as f:
            full_tex = f.read()
    except Exception:
        return False

    def _strip_tex(s: str) -> str:
        t = str(s or "")
        t = re.sub(r"%.*$", "", t, flags=re.MULTILINE)
        t = re.sub(r"\\(textbf|textit|emph|underline|texttt)\{([^\}]*)\}", r"\2", t)
        t = re.sub(r"\\(section|subsection)\*?\{([^\}]*)\}", r"\2", t)
        t = re.sub(r"\\(vspace|hspace|textwidth|linewidth|centering|column|footnotesize|small|large|Huge)", "", t)
        t = re.sub(r"\\(begin|end)\{(itemize|enumerate|columns|column|tabular|center|figure|table)\}(\[[^\]]*\])?", "", t)
        t = re.sub(r"\\item\s", "• ", t)
        t = re.sub(r"\\item\[([^\]]*)\]", r"\1 ", t)
        t = re.sub(r"\\([a-zA-Z]+)(\[[^\]]*\])?(\{([^\}]*)\})?", " ", t)
        t = re.sub(r"[\{\}]", "", t)
        t = re.sub(r"\s+", " ", t)
        return t.strip()

    def _find_images(frame_tex: str) -> List[str]:
        imgs = []
        for im in re.findall(r"\\includegraphics\*?(?:\[[^\]]*\])?\{([^\}]+)\}", frame_tex):
            rel = str(im).strip()
            candidates = [
                os.path.join(recipe_dir, rel),
                os.path.join(os.path.dirname(recipe_dir), rel),
                os.path.join(recipe_dir, rel + ".png"),
                os.path.join(recipe_dir, rel + ".jpg"),
            ]
            for c in candidates:
                if os.path.exists(c) and os.path.isfile(c):
                    imgs.append(c)
                    break
        return imgs

    def _extract_title(frame_tex: str) -> str:
        m = re.search(r"\\frametitle\{([^}]*)\}", frame_tex)
        if m:
            return _strip_tex(m.group(1))
        return ""

    def _extract_bullets(frame_tex: str) -> List[str]:
        body = re.sub(r"\\frametitle\{[^\}]*\}", "", frame_tex)
        body = re.sub(r"\\framesubtitle\{[^\}]*\}", "", body)
        body = re.sub(r"\\includegraphics[^\}]*\}", "", body)

        items = []
        lines = body.splitlines()
        current_item = []

        for ln in lines:
            ln = ln.strip()
            if not ln or ln.startswith("%"):
                continue
            if "\\item" in ln:
                if current_item:
                    items.append(_strip_tex(" ".join(current_item)))
                current_item = [ln.replace("\\item", "")]
            else:
                if current_item:
                    current_item.append(ln)
                elif ln and not ln.startswith("\\"):
                    pass
        if current_item:
            items.append(_strip_tex(" ".join(current_item)))
        return items

    raw_frames = re.findall(r"\\begin\{frame\}(.*?)\\end\{frame\}", full_tex, re.DOTALL)
    if not raw_frames:
        raw_frames = re.findall(r"\\frame\{(.*?)\}", full_tex, re.DOTALL)

    if not raw_frames:
        return False

    prs = Presentation()
    prs.slide_width = Emu(int(13.333 * 914400))
    prs.slide_height = Emu(int(7.5 * 914400))

    layout_content = prs.slide_layouts[1]
    slides_created = 0

    for fr in raw_frames:
        title = _extract_title(fr)
        bullets = _extract_bullets(fr)
        images = _find_images(fr)

        slide = prs.slides.add_slide(layout_content)

        try:
            if slide.shapes.title:
                slide.shapes.title.text = title
        except Exception:
            pass

        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break

        if body_shape and body_shape.has_text_frame:
            tf = body_shape.text_frame
            tf.clear()

            if bullets:
                for i, b in enumerate(bullets):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = b
                    p.level = 0
                    p.space_after = Pt(10)
            elif not images:
                plain = _strip_tex(re.sub(r"\\frametitle\{[^\}]*\}", "", fr))
                if plain:
                    p = tf.paragraphs[0]
                    p.text = plain

        if images:
            img_path = images[0]
            sw = prs.slide_width
            sh = prs.slide_height

            try:
                if bullets:
                    left = Emu(sw * 0.55)
                    top = Emu(sh * 0.2)
                    height = Emu(sh * 0.6)
                    slide.shapes.add_picture(img_path, left, top, height=height)

                    if body_shape:
                        body_shape.width = Emu(sw * 0.5)
                else:
                    height = Emu(sh * 0.7)
                    pic = slide.shapes.add_picture(img_path, 0, 0, height=height)
                    pic.left = int((sw - pic.width) / 2)
                    pic.top = int((sh - pic.height) / 2)
            except Exception:
                pass

        slides_created += 1

    if slides_created == 0:
        return False

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return True


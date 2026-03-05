from __future__ import annotations

from dataclasses import asdict
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation

from ..io_utils import sha1_bytes
from .common import DeckContent, ExtractorResult, SlideText


def _iter_text_runs(slide) -> Tuple[List[str], List[float]]:
    texts: List[str] = []
    font_pts: List[float] = []
    for shape in slide.shapes:
        if not hasattr(shape, "has_text_frame"):
            continue
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                t = (run.text or "").strip()
                if t:
                    texts.append(t)
                if run.font is not None and run.font.size is not None:
                    try:
                        font_pts.append(float(run.font.size.pt))
                    except Exception:
                        pass
    return texts, font_pts


def _extract_notes(slide) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    if notes_slide is None:
        return ""
    texts: List[str] = []
    for shape in notes_slide.shapes:
        if not hasattr(shape, "has_text_frame"):
            continue
        if not shape.has_text_frame:
            continue
        t = (shape.text or "").strip()
        if t:
            texts.append(t)
    return "\n".join(texts).strip()


def extract_deck_from_pptx(path: Path) -> ExtractorResult:
    diagnostics = {"path": str(path)}
    try:
        pres = Presentation(str(path))
    except Exception as e:
        return ExtractorResult(ok=False, error=f"open_pptx_failed: {e}", content=None, diagnostics=diagnostics)

    slides_out: List[SlideText] = []
    deck_text_parts: List[str] = []
    deck_notes_parts: List[str] = []
    image_hashes: List[str] = []

    for idx, slide in enumerate(pres.slides):
        text_runs, font_pts = _iter_text_runs(slide)
        slide_text = "\n".join(text_runs).strip()
        notes_text = _extract_notes(slide)

        num_images = 0
        for shape in slide.shapes:
            try:
                if shape.shape_type == 13:
                    num_images += 1
                    image_hashes.append(sha1_bytes(shape.image.blob))
            except Exception:
                continue

        min_font = min(font_pts) if font_pts else None
        words = len([w for w in slide_text.replace("\n", " ").split(" ") if w])

        slides_out.append(
            SlideText(
                slide_index=idx,
                text=slide_text,
                text_ocr="",
                text_final=slide_text,
                notes=notes_text,
                word_count=words,
                min_font_pt=min_font,
                num_shapes=len(slide.shapes),
                num_images=num_images,
                ocr_used=False,
                ocr_confidence=None,
            )
        )

        if slide_text:
            deck_text_parts.append(slide_text)
        if notes_text:
            deck_notes_parts.append(notes_text)

    deck_text = "\n\n".join(deck_text_parts).strip()
    content = DeckContent(
        artifact_path=str(path),
        artifact_type="pptx",
        slides=slides_out,
        deck_text=deck_text,
        deck_text_final=deck_text,
        deck_notes="\n\n".join(deck_notes_parts).strip(),
        image_hashes=image_hashes,
    )
    diagnostics["num_slides"] = len(slides_out)
    return ExtractorResult(ok=True, error=None, content=content, diagnostics=diagnostics)

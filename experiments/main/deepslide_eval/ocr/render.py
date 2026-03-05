from __future__ import annotations

import subprocess
from pathlib import Path
from typing import List

from PIL import Image


def render_pdf_to_images(pdf_path: Path, out_dir: Path, dpi: int = 200) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        import fitz
    except Exception as e:
        raise RuntimeError(f"missing_pymupdf: {e}")

    doc = fitz.open(str(pdf_path))
    out: List[Path] = []
    try:
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            p = out_dir / f"slide_{i:03d}.png"
            pix.save(str(p))
            out.append(p)
    finally:
        doc.close()
    return out


def convert_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx_path),
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    out_text = (proc.stdout or "") + "\n" + (proc.stderr or "")
    if proc.returncode != 0 or "Error:" in out_text:
        raise RuntimeError(f"pptx_to_pdf_failed: {out_text.strip()}")
    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        pdfs = list(out_dir.glob("*.pdf"))
        if pdfs:
            pdfs = sorted(pdfs, key=lambda p: p.stat().st_mtime, reverse=True)
            return pdfs[0]
        raise FileNotFoundError(str(pdf_path))
    return pdf_path


def extract_pptx_slide_images(pptx_path: Path, out_dir: Path) -> List[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError(f"missing_python_pptx: {e}")

    prs = Presentation(str(pptx_path))
    out: List[Path] = []
    for idx, slide in enumerate(prs.slides):
        best = None
        best_area = 0
        for shape in slide.shapes:
            if not getattr(shape, "shape_type", None):
                continue
            if not getattr(shape, "image", None):
                continue
            try:
                area = int(shape.width) * int(shape.height)
            except Exception:
                area = 0
            if area >= best_area:
                best_area = area
                best = shape

        if best is None:
            out.append(out_dir / f"slide_{idx:03d}.png")
            continue

        blob = best.image.blob
        ext = (best.image.ext or "png").lower()
        p = out_dir / f"slide_{idx:03d}.{ext}"
        p.write_bytes(blob)
        if ext not in {"png", "jpg", "jpeg"}:
            try:
                with Image.open(p) as im:
                    p2 = out_dir / f"slide_{idx:03d}.png"
                    im.save(p2)
                p = p2
            except Exception:
                pass
        out.append(p)
    return out

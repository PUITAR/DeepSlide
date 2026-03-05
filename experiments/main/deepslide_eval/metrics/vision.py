from __future__ import annotations

import io
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

import numpy as np


@dataclass(frozen=True)
class ImageItem:
    key: str
    width: int
    height: int
    data: bytes


def _sha1_bytes(data: bytes) -> str:
    import hashlib

    h = hashlib.sha1()
    h.update(data)
    return h.hexdigest()


def extract_images_from_pdf(pdf_path: Path) -> Tuple[List[ImageItem], Dict]:
    diagnostics: Dict = {"pdf_path": str(pdf_path)}
    try:
        import fitz
    except Exception as e:
        return [], {"ok": False, "reason": f"missing_pymupdf: {e}", **diagnostics}

    items: List[ImageItem] = []
    try:
        doc = fitz.open(str(pdf_path))
    except Exception as e:
        return [], {"ok": False, "reason": f"open_pdf_failed: {e}", **diagnostics}

    try:
        for page in doc:
            try:
                imgs = page.get_images(full=True)
            except Exception:
                imgs = []
            for img in imgs:
                try:
                    xref = img[0]
                    base = doc.extract_image(xref)
                    b = base.get("image")
                    if not isinstance(b, (bytes, bytearray)):
                        continue
                    b = bytes(b)
                    if not b:
                        continue
                    w = int(base.get("width") or 0)
                    h = int(base.get("height") or 0)
                    key = _sha1_bytes(b)
                    items.append(ImageItem(key=key, width=w, height=h, data=b))
                except Exception:
                    continue
    finally:
        doc.close()

    diagnostics["ok"] = True
    diagnostics["num_images_raw"] = len(items)
    return items, diagnostics


def extract_images_from_pptx(pptx_path: Path) -> Tuple[List[ImageItem], Dict]:
    diagnostics: Dict = {"pptx_path": str(pptx_path)}
    try:
        from pptx import Presentation
    except Exception as e:
        return [], {"ok": False, "reason": f"missing_python_pptx: {e}", **diagnostics}

    try:
        pres = Presentation(str(pptx_path))
    except Exception as e:
        return [], {"ok": False, "reason": f"open_pptx_failed: {e}", **diagnostics}

    items: List[ImageItem] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            try:
                if shape.shape_type != 13:
                    continue
                b = shape.image.blob
                if not isinstance(b, (bytes, bytearray)):
                    continue
                b = bytes(b)
                if not b:
                    continue
                w = 0
                h = 0
                try:
                    sz = getattr(shape.image, "size", None)
                    if sz is not None and len(sz) == 2:
                        w = int(sz[0])
                        h = int(sz[1])
                except Exception:
                    pass
                key = _sha1_bytes(b)
                items.append(ImageItem(key=key, width=w, height=h, data=b))
            except Exception:
                continue

    diagnostics["ok"] = True
    diagnostics["num_images_raw"] = len(items)
    return items, diagnostics


def filter_images(
    items: Sequence[ImageItem],
    min_edge_px: int = 64,
    min_area_px: int = 64 * 64,
) -> Tuple[List[ImageItem], Dict]:
    from PIL import Image

    kept: List[ImageItem] = []
    rejected_small = 0
    rejected_decode = 0
    for it in items:
        w = it.width
        h = it.height
        if w <= 0 or h <= 0:
            try:
                im = Image.open(io.BytesIO(it.data))
                w, h = im.size
            except Exception:
                rejected_decode += 1
                continue
        if min(w, h) < min_edge_px or (w * h) < min_area_px:
            rejected_small += 1
            continue
        kept.append(ImageItem(key=it.key, width=int(w), height=int(h), data=it.data))
    return kept, {"kept": len(kept), "rejected_small": rejected_small, "rejected_decode": rejected_decode}


class VisionEmbedder:
    def __init__(self) -> None:
        self._device = None
        self._processor = None
        self._model = None
        self._cache: Dict[str, np.ndarray] = {}

    @property
    def model_id(self) -> str:
        return (os.getenv("EVAL_FVIS_VIT_MODEL") or "google/vit-base-patch16-224-in21k").strip()

    @property
    def device(self):
        if self._device is None:
            import torch

            env = (os.getenv("EVAL_FVIS_DEVICE") or "").strip()
            if env:
                self._device = torch.device(env)
            else:
                self._device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        return self._device

    def _ensure_loaded(self) -> None:
        if self._processor is not None and self._model is not None:
            return
        import torch
        from transformers import AutoImageProcessor, AutoModel

        self._processor = AutoImageProcessor.from_pretrained(self.model_id)
        self._model = AutoModel.from_pretrained(self.model_id)
        self._model.eval()
        self._model.to(self.device)
        if self.device.type == "cuda":
            self._model.to(dtype=torch.float16)

    def embed_images(self, items: Sequence[ImageItem], batch_size: int = 24) -> np.ndarray:
        from PIL import Image
        import torch

        self._ensure_loaded()
        assert self._processor is not None
        assert self._model is not None

        vecs: List[np.ndarray] = []
        for i in range(0, len(items), batch_size):
            batch = items[i : i + batch_size]
            imgs = []
            keys = []
            for it in batch:
                if it.key in self._cache:
                    continue
                try:
                    im = Image.open(io.BytesIO(it.data)).convert("RGB")
                except Exception:
                    continue
                imgs.append(im)
                keys.append(it.key)
            if imgs:
                inputs = self._processor(images=imgs, return_tensors="pt")
                inputs = {k: v.to(self.device) for k, v in inputs.items()}
                with torch.no_grad():
                    out = self._model(**inputs)
                    h = out.last_hidden_state[:, 0, :]
                    h = torch.nn.functional.normalize(h, p=2, dim=1)
                    h = h.detach().float().cpu().numpy().astype(np.float32)
                for k, v in zip(keys, h):
                    self._cache[k] = v

            for it in batch:
                v = self._cache.get(it.key)
                if v is None:
                    v = np.zeros((768,), dtype=np.float32)
                    v[0] = np.nan
                vecs.append(v)

        return np.stack(vecs, axis=0) if vecs else np.zeros((0, 768), dtype=np.float32)


def compute_f_vis_vit(
    source_images: Sequence[ImageItem],
    deck_images: Sequence[ImageItem],
    embedder: VisionEmbedder,
    threshold: float,
) -> Tuple[Optional[float], Dict]:
    if not source_images:
        return None, {"reason": "no_source_images"}
    if not deck_images:
        return None, {"reason": "no_deck_images"}

    src_emb = embedder.embed_images(source_images)
    out_emb = embedder.embed_images(deck_images)
    if src_emb.shape[0] == 0:
        return None, {"reason": "no_source_embeddings"}
    if out_emb.shape[0] == 0:
        return None, {"reason": "no_deck_embeddings"}

    src_ok = ~np.isnan(src_emb[:, 0])
    out_ok = ~np.isnan(out_emb[:, 0])
    src_idx = np.where(src_ok)[0]
    out_idx = np.where(out_ok)[0]
    if src_idx.size == 0:
        return None, {"reason": "no_source_embeddings"}
    if out_idx.size == 0:
        return None, {"reason": "no_deck_embeddings"}

    src = src_emb[src_idx]
    out = out_emb[out_idx]

    sim = np.clip(src @ out.T, -1.0, 1.0).astype(np.float32)
    n = sim.shape[0]
    m = sim.shape[1]
    big = 1e6
    cost = np.full((n + m, n + m), 0.0, dtype=np.float64)
    sub = np.where(sim >= float(threshold), -sim.astype(np.float64), big)
    cost[:n, :m] = sub
    cost[:n, m:] = 0.0
    cost[n:, :m] = 0.0

    from scipy.optimize import linear_sum_assignment

    r, c = linear_sum_assignment(cost)
    matched_pairs = []
    for rr, cc in zip(r.tolist(), c.tolist()):
        if rr < n and cc < m and cost[rr, cc] < big / 10:
            matched_pairs.append((int(src_idx[rr]), int(out_idx[cc]), float(sim[rr, cc])))

    matched = len(matched_pairs)
    score = float(matched / max(1, len(source_images)))
    mean_sim = float(np.mean([s for _, _, s in matched_pairs])) if matched_pairs else 0.0

    return score, {
        "method": "vit_cls_cosine",
        "model": embedder.model_id,
        "threshold": float(threshold),
        "src": int(len(source_images)),
        "out": int(len(deck_images)),
        "matched": int(matched),
        "mean_sim_matched": mean_sim,
    }

